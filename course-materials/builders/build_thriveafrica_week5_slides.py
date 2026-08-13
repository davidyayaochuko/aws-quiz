from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = r"C:\Users\Yaya\Projects\thriveafrica-aws-quiz\course-materials\week-5-storage\ThriveAfrica_Week5_Storage_Class.pptx"

NAVY = RGBColor(12, 26, 43)
NAVY_2 = RGBColor(22, 44, 68)
ORANGE = RGBColor(255, 153, 0)
BLUE = RGBColor(64, 155, 245)
GREEN = RGBColor(66, 193, 123)
PURPLE = RGBColor(161, 108, 244)
RED = RGBColor(235, 91, 104)
WHITE = RGBColor(249, 251, 253)
PALE = RGBColor(240, 244, 248)
INK = RGBColor(31, 45, 59)
MUTED = RGBColor(104, 120, 136)
LIGHT_LINE = RGBColor(216, 225, 233)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, radius=False, line=None):
    kind = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_text(
    slide,
    value,
    x,
    y,
    w,
    h,
    size=22,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, size=20, color=INK, gap=9):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = "•  " + item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
    return box


def base(title, subtitle="", section="WEEK 5 • STORAGE"):
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PALE
    add_rect(slide, 0, 0, 13.333, 0.14, ORANGE)
    add_text(slide, section, 0.65, 0.38, 4.5, 0.3, 11, ORANGE, True)
    add_text(slide, title, 0.65, 0.77, 12.0, 0.7, 30, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 0.67, 1.43, 11.8, 0.4, 15, MUTED)
    add_text(slide, "ThriveAfrica • AWS Cloud Practitioner", 0.65, 7.08, 5.0, 0.2, 10, MUTED)
    add_text(slide, str(len(prs.slides)), 12.15, 7.05, 0.5, 0.2, 10, MUTED, align=PP_ALIGN.RIGHT)
    return slide


def card(slide, x, y, w, h, heading, body, accent=ORANGE, body_size=16):
    add_rect(slide, x, y, w, h, WHITE, True, LIGHT_LINE)
    add_rect(slide, x, y, 0.09, h, accent)
    add_text(slide, heading, x + 0.25, y + 0.2, w - 0.45, 0.42, 18, NAVY, True)
    add_text(slide, body, x + 0.25, y + 0.7, w - 0.45, h - 0.82, body_size, MUTED)


def section(title, subtitle, number):
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_rect(slide, 0, 0, 0.18, 7.5, ORANGE)
    add_text(slide, f"0{number}", 0.85, 1.08, 1.3, 0.8, 36, ORANGE, True)
    add_text(slide, title, 0.85, 2.0, 11.5, 1.25, 41, WHITE, True)
    add_text(slide, subtitle, 0.87, 3.58, 10.8, 1.0, 22, RGBColor(180, 194, 208))
    add_text(slide, "THRIVEAFRICA • WEEK 5", 0.87, 6.65, 4.5, 0.3, 12, WHITE, True)
    return slide


def question_slide(question, options, prompt="Vote, then explain your choice."):
    slide = base("Ask the room", prompt)
    add_text(slide, question, 0.8, 2.0, 11.7, 0.85, 24, NAVY, True)
    colors = [BLUE, GREEN, ORANGE, PURPLE]
    for i, option in enumerate(options):
        row, col = divmod(i, 2)
        x = 0.85 + col * 6.15
        y = 3.15 + row * 1.22
        add_rect(slide, x, y, 5.65, 0.9, WHITE, True, LIGHT_LINE)
        add_rect(slide, x, y, 0.68, 0.9, colors[i], True, colors[i])
        add_text(slide, chr(65 + i), x + 0.11, y + 0.22, 0.44, 0.3, 18, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, option, x + 0.88, y + 0.2, 4.5, 0.45, 17, INK)
    add_text(slide, "Answer is in the instructor guide.", 0.85, 5.85, 11.7, 0.35, 14, MUTED, align=PP_ALIGN.CENTER)
    return slide


# 1 — Title
slide = prs.slides.add_slide(blank)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY
add_rect(slide, 0, 0, 13.333, 0.18, ORANGE)
add_text(slide, "AWS CLOUD PRACTITIONER • WEEK 5", 0.82, 1.0, 6.5, 0.4, 16, ORANGE, True)
add_text(slide, "Storage on AWS", 0.82, 1.65, 8.5, 0.8, 43, WHITE, True)
add_text(slide, "S3 • EBS • EFS • FSx • Gateway • Backup", 0.85, 2.75, 9.5, 0.5, 22, RGBColor(186, 200, 214))
add_rect(slide, 9.55, 1.3, 2.25, 2.25, ORANGE, True)
add_text(slide, "05", 9.55, 1.85, 2.25, 0.85, 39, NAVY, True, PP_ALIGN.CENTER)
add_text(slide, "ThriveAfrica", 0.85, 6.65, 4.0, 0.3, 14, WHITE, True)

# 2 — Objectives
slide = base("Today’s destination", "Choose the right storage service from a workload scenario.")
items = [
    ("TYPES", "Separate block, file, and object storage", ORANGE),
    ("S3", "Buckets, classes, lifecycle, protection", BLUE),
    ("BLOCK", "EBS volumes, snapshots, instance store", GREEN),
    ("HYBRID", "EFS, FSx, Gateway, Backup, DRS", PURPLE),
]
for i, (head, body, color) in enumerate(items):
    x = 0.8 + i * 3.1
    card(slide, x, 2.3, 2.75, 2.7, head, body, color, 16)
add_text(slide, "Goal: identify the constraint → eliminate poor fits → select the best storage option.", 0.85, 5.55, 11.7, 0.5, 18, NAVY, True, PP_ALIGN.CENTER)

# 3 — Opening scenario
question_slide(
    "A company has website images, a database disk, and a shared Linux folder. Same storage type?",
    ["All object storage", "All block storage", "Three different models", "Only file storage"],
    "Name the model that fits each workload.",
)

# 4
section("Three storage models", "Block, file, and object answer different exam constraints.", 1)

# 5
slide = base("Block • File • Object", "Match the mental model before naming a service.")
card(slide, 0.75, 2.0, 3.85, 3.55, "Block", "Virtual hard disk\n\nOne server’s volume\nBoot disks • databases\n\nEBS • instance store", ORANGE, 17)
card(slide, 4.75, 2.0, 3.85, 3.55, "File", "Shared folder paths\n\nMany servers mount it\nNFS or SMB shares\n\nEFS • FSx", BLUE, 17)
card(slide, 8.75, 2.0, 3.85, 3.55, "Object", "Objects in a bucket\n\nAccessed by key\nWeb assets • archives\n\nAmazon S3", GREEN, 17)

# 6
question_slide(
    "Static product images for a public website need durable object storage.",
    ["Amazon EBS", "Amazon S3", "Instance store", "Amazon FSx for Windows"],
)

# 7
section("Amazon S3 foundations", "Object storage with durability, control, and lifecycle automation.", 2)

# 8
slide = base("Buckets and objects", "S3 stores objects inside regionally created buckets.")
add_bullets(
    slide,
    [
        "Bucket names must be globally unique",
        "Choose a Region when you create the bucket",
        "An object has a key, data, and metadata",
        "Access via console, CLI, SDK, or API over HTTPS",
        "Ideal for backups, static assets, data lakes, archives",
    ],
    0.8,
    2.05,
    6.2,
    3.8,
    19,
)
card(slide, 7.35, 2.1, 5.0, 1.5, "Remember", "S3 = object storage, not a block disk", ORANGE, 17)
card(slide, 7.35, 3.9, 5.0, 1.7, "Exam cue", "Unlimited objects • URL-style keys • regional buckets", BLUE, 17)

# 9
slide = base("Durability vs availability", "Two words the exam loves to swap.")
card(slide, 0.85, 2.05, 5.55, 3.4, "Durability", "Will the data survive without loss?\n\nS3 Standard is designed for\n11 nines durability", ORANGE, 18)
card(slide, 7.0, 2.05, 5.55, 3.4, "Availability", "Can you retrieve it when needed?\n\nClasses differ in retrieval\nspeed and availability design", BLUE, 18)
add_text(slide, "Durability = do not lose it     •     Availability = can I get it now?", 0.85, 5.8, 11.7, 0.4, 18, NAVY, True, PP_ALIGN.CENTER)

# 10
slide = base("Protect and control S3 data", "Versioning, Block Public Access, and Object Lock.")
card(slide, 0.75, 2.0, 3.85, 3.55, "Versioning", "Keep prior versions\nRecover overwrites\nDelete markers\n\nRestore earlier copies", ORANGE, 17)
card(slide, 4.75, 2.0, 3.85, 3.55, "Block Public Access", "Stop accidental public\nbucket or object access\n\nDefault on new buckets", BLUE, 17)
card(slide, 8.75, 2.0, 3.85, 3.55, "Object Lock", "WORM retention\nCompliance immutability\n\nRequires versioning", GREEN, 17)

# 11
slide = base("Lifecycle and replication", "Automate cost tiering; copy objects when required.")
card(slide, 0.85, 2.05, 5.55, 3.5, "Lifecycle policies", "Transition or expire automatically\n\nLab pattern:\nIA at 30 days\nGlacier at 90 days\nExpire at 365 days", ORANGE, 17)
card(slide, 7.0, 2.05, 5.55, 3.5, "Replication", "Copy to another bucket\nSRR or Cross-Region (CRR)\n\nVersioning required\nDR / compliance clue → CRR", BLUE, 17)

# 12
section("Seven storage classes", "Access pattern + retrieval time + minimum duration.", 3)

# 13
slide = base("The seven S3 storage classes", "Know the job of each class.")
classes = [
    ("Standard", "Frequent access", ORANGE),
    ("Intelligent-Tiering", "Unknown patterns", BLUE),
    ("Standard-IA", "Infrequent + fast", GREEN),
    ("One Zone-IA", "Infrequent + one AZ", PURPLE),
    ("Glacier Instant", "Archive + ms retrieve", RED),
    ("Glacier Flexible", "Archive + hours OK", ORANGE),
    ("Deep Archive", "Lowest-cost archive", BLUE),
]
for i, (head, body, color) in enumerate(classes):
    x = 0.55 + i * 1.82
    add_rect(slide, x, 2.15, 1.7, 3.35, WHITE, True, LIGHT_LINE)
    add_rect(slide, x, 2.15, 1.7, 0.12, color)
    add_text(slide, head, x + 0.08, 2.5, 1.54, 1.1, 14, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, body, x + 0.08, 3.9, 1.54, 1.1, 13, MUTED, align=PP_ALIGN.CENTER)

# 14
slide = base("Retrieval times and minimum durations", "These numbers win scenario questions.")
rows = [
    ("Standard", "Milliseconds", "None", ORANGE),
    ("Intelligent-Tiering", "ms (freq/infreq)", "Watch archive tiers", BLUE),
    ("Standard-IA", "Milliseconds", "30 days", GREEN),
    ("One Zone-IA", "Milliseconds", "30 days", PURPLE),
    ("Glacier Instant Retrieval", "Milliseconds", "90 days", RED),
    ("Glacier Flexible Retrieval", "Minutes to hours", "90 days", ORANGE),
    ("Glacier Deep Archive", "~12h / up to ~48h", "180 days", BLUE),
]
add_rect(slide, 0.75, 1.9, 11.85, 0.5, NAVY_2, True)
add_text(slide, "Class", 0.95, 2.0, 3.8, 0.3, 14, WHITE, True)
add_text(slide, "Retrieval", 5.0, 2.0, 3.5, 0.3, 14, WHITE, True)
add_text(slide, "Minimum duration", 8.7, 2.0, 3.5, 0.3, 14, WHITE, True)
for i, (cls, retrieval, minimum, color) in enumerate(rows):
    y = 2.5 + i * 0.55
    add_rect(slide, 0.75, y, 11.85, 0.48, WHITE, True, LIGHT_LINE)
    add_rect(slide, 0.75, y, 0.1, 0.48, color)
    add_text(slide, cls, 1.0, y + 0.08, 3.8, 0.3, 13, NAVY, True)
    add_text(slide, retrieval, 5.0, y + 0.08, 3.5, 0.3, 13, INK)
    add_text(slide, minimum, 8.7, y + 0.08, 3.5, 0.3, 13, MUTED)

# 15
question_slide(
    "Data must be archived at lowest cost. Retrieval in about 12 hours is acceptable.",
    ["S3 Standard", "S3 Glacier Instant Retrieval", "S3 Glacier Deep Archive", "S3 One Zone-IA"],
)

# 16
section("Block storage for EC2", "Persistent EBS volumes versus temporary instance store.", 4)

# 17
slide = base("EBS volume types", "Match the performance clue, not every IOPS number.")
types = [
    ("gp3 / gp2", "General-purpose SSD\nBoot volumes\nEveryday workloads", ORANGE),
    ("io1 / io2", "Provisioned IOPS SSD\nConsistent high IOPS\nCritical databases", BLUE),
    ("st1", "Throughput HDD\nLarge sequential data\nBig data / logs", GREEN),
    ("sc1", "Cold HDD\nInfrequent access\nLowest-cost HDD", PURPLE),
]
for i, (head, body, color) in enumerate(types):
    x = 0.75 + i * 3.15
    card(slide, x, 2.1, 2.95, 3.4, head, body, color, 16)

# 18
slide = base("Snapshots and Availability Zones", "An EBS volume lives in one AZ.")
add_bullets(
    slide,
    [
        "Create a snapshot of the volume",
        "Snapshots are stored as AWS-managed backups in S3",
        "Snapshots are incremental after the first full copy",
        "Create a new volume from the snapshot in another AZ",
        "You cannot simply re-attach the same volume across AZs",
    ],
    0.85,
    2.05,
    7.0,
    3.6,
    18,
)
card(slide, 8.2, 2.2, 4.2, 3.2, "Lab move", "AZ A volume\n→ snapshot\n→ AZ B volume", ORANGE, 18)

# 19
slide = base("Instance store is temporary", "Local host disk is fast—and ephemeral.")
card(slide, 0.85, 2.05, 5.55, 3.45, "Use instance store for", "Caches\nBuffers\nScratch space\nTemporary processing", BLUE, 18)
card(slide, 7.0, 2.05, 5.55, 3.45, "Do not use it for", "Only copy of important data\nDatabase durability\nAnything that must survive stop/terminate", RED, 18)
add_text(slide, "EBS persists  •  Instance store is lost on stop, terminate, or host failure", 0.85, 5.85, 11.7, 0.4, 17, NAVY, True, PP_ALIGN.CENTER)

# 20
section("File systems and hybrid storage", "Shared files in AWS, plus on-premises bridges.", 5)

# 21
slide = base("EFS vs FSx for Windows File Server", "Protocol and operating system decide.")
card(slide, 0.85, 2.05, 5.55, 3.5, "Amazon EFS", "NFS file system\nLinux / Unix workloads\nRegional multi-AZ design\n\nMany EC2 mounts across AZs", ORANGE, 17)
card(slide, 7.0, 2.05, 5.55, 3.5, "FSx for Windows", "SMB file shares\nWindows workloads\nActive Directory integration\n\nManaged Windows file server", BLUE, 17)

# 22
slide = base("AWS Storage Gateway", "Connect on-premises environments to AWS storage.")
card(slide, 0.75, 2.05, 3.85, 3.45, "File Gateway", "NFS / SMB shares\nFiles stored as\nS3 objects", ORANGE, 17)
card(slide, 4.75, 2.05, 3.85, 3.45, "Volume Gateway", "iSCSI volumes\nCached or stored\nvolume modes", BLUE, 17)
card(slide, 8.75, 2.05, 3.85, 3.45, "Tape Gateway", "Virtual tape library\nReplace physical\ntape workflows", GREEN, 17)

# 23
question_slide(
    "On-premises servers need SMB/NFS shares that store files in Amazon S3.",
    ["Tape Gateway", "File Gateway", "Instance store", "Gateway Load Balancer"],
)

# 24
slide = base("AWS Backup and Elastic Disaster Recovery", "Protect copies versus recover whole systems.")
card(slide, 0.85, 2.05, 5.55, 3.5, "AWS Backup", "Central backup plans\nRetention policies\nAcross supported services\n\nClue: manage backups in one place", ORANGE, 17)
card(slide, 7.0, 2.05, 5.55, 3.5, "Elastic Disaster Recovery", "Replicate servers\nFail over into AWS\nReduce downtime after disruption\n\nClue: recover systems fast", BLUE, 17)

# 25
question_slide(
    "A company wants one service to schedule and retain backups across EBS, EFS, and databases.",
    ["Amazon CloudFront", "AWS Backup", "Amazon Route 53", "AWS WAF"],
)

# 26
slide = base("Service selection sprint", "Call out the service before checking the answer key.")
scenarios = [
    ("Prevent accidental public S3 access", "Block Public Access"),
    ("Unknown object access pattern", "Intelligent-Tiering"),
    ("Shared Linux file system across AZs", "Amazon EFS"),
    ("Windows SMB + Active Directory", "FSx for Windows"),
    ("Replace physical tapes", "Tape Gateway"),
]
for i, (scenario, answer) in enumerate(scenarios):
    y = 1.92 + i * 0.84
    add_rect(slide, 0.8, y, 7.65, 0.64, WHITE, True, LIGHT_LINE)
    add_text(slide, scenario, 1.05, y + 0.15, 7.1, 0.3, 16, NAVY, True)
    add_rect(slide, 8.75, y, 3.6, 0.64, NAVY_2, True)
    add_text(slide, answer, 8.75, y + 0.15, 3.6, 0.3, 15, WHITE, True, PP_ALIGN.CENTER)

# 27
slide = base("Five exam traps", "Know what each option does—and does not do.")
traps = [
    "Durability is not the same as availability",
    "One Zone-IA is single-AZ, not multi-AZ",
    "Glacier Flexible / Deep Archive are not millisecond retrieval",
    "EBS volumes do not span Availability Zones; use snapshots to rebuild",
    "Instance store is temporary; EFS is not the Windows SMB answer",
]
add_bullets(slide, traps, 0.9, 2.05, 11.5, 3.9, 20, NAVY, 14)

# 28
slide = base("Homework lab preview", "Version it, tier it, then take it down.")
steps = [
    ("1", "Create bucket + confirm Block Public Access"),
    ("2", "Enable versioning; overwrite and restore"),
    ("3", "Lifecycle: IA 30 → Glacier 90 → expire 365"),
    ("4", "Change one object storage class manually"),
    ("5", "EBS volume → snapshot → rebuild in another AZ"),
    ("6", "Delete objects, bucket, volumes, and snapshot"),
]
for i, (num, text) in enumerate(steps):
    row, col = divmod(i, 2)
    x = 0.8 + col * 6.15
    y = 1.95 + row * 1.35
    add_rect(slide, x, y, 5.9, 1.1, WHITE, True, LIGHT_LINE)
    add_rect(slide, x, y, 0.7, 1.1, ORANGE if col == 0 else BLUE, True)
    add_text(slide, num, x, y + 0.32, 0.7, 0.35, 20, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, text, x + 0.9, y + 0.32, 4.7, 0.5, 15, NAVY, True)

# 29
slide = base("Cost and teardown warning", "Small labs still create billable leftovers.")
add_bullets(
    slide,
    [
        "Empty all object versions before deleting a bucket",
        "Delete EBS volumes and snapshots after the rebuild test",
        "IA and Glacier classes can have minimum-duration charges",
        "Confirm budget alerts from earlier weeks are still active",
        "Verify the console shows no unexpected lab resources",
    ],
    0.9,
    2.05,
    11.5,
    3.8,
    20,
)

# 30
slide = prs.slides.add_slide(blank)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY
add_rect(slide, 0, 0, 13.333, 0.18, ORANGE)
add_text(slide, "WEEK 5 RECAP", 0.85, 0.72, 4.0, 0.35, 15, ORANGE, True)
add_text(slide, "Choose storage from the constraint", 0.85, 1.28, 11.7, 0.7, 33, WHITE, True)
recap = [
    ("MODELS", "Block disk • shared file • object bucket"),
    ("S3", "Classes, lifecycle, versioning, BPA, Object Lock"),
    ("BLOCK", "EBS persists; snapshots move AZs; instance store is temporary"),
    ("HYBRID", "EFS/FSx, Storage Gateway, Backup, Elastic Disaster Recovery"),
]
for i, (head, body) in enumerate(recap):
    y = 2.35 + i * 0.87
    add_rect(slide, 0.9, y, 11.5, 0.66, NAVY_2, True)
    add_text(slide, head, 1.2, y + 0.15, 1.7, 0.3, 16, ORANGE, True)
    add_text(slide, body, 3.1, y + 0.15, 8.9, 0.3, 16, WHITE)
add_text(slide, "Complete the Week 5 knowledge check (14/20) and lab teardown.", 0.9, 6.45, 11.5, 0.4, 17, RGBColor(191, 204, 217), True, PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Created {OUT} with {len(prs.slides)} slides")
